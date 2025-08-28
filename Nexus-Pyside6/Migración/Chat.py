import sys
import os
import json
import httpx
from pptx import Presentation
from PySide6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QTextEdit, QPushButton, QLabel, QMessageBox, QProgressBar, QFrame,
    QPlainTextEdit, QListWidget, QListWidgetItem, QStackedWidget
)
from PySide6.QtCore import Qt, Signal, QRunnable, QThreadPool, QObject, QEvent, QSize, Qt

# --- Lógica de la aplicación principal ---

# La clave de API se inyecta automáticamente en el entorno de Canvas.
API_KEY = "AIzaSyCrNYH7OtSt7c9uxkSJ9LE1s0YnFSE-e9U"
API_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={API_KEY}"
PLAN_CAPACITACION_FILE = r"D:\Proyectos_python\PyCharmMiscProject\PLAN de Capacitacion.pptx"


def cargar_conocimiento(path):
    """
    Carga el texto de todas las formas de un archivo PowerPoint.
    """
    # print(f"🔄 Intentando cargar el archivo PowerPoint: {path}")
    texto = ""
    try:
        if not os.path.exists(path):
            return "❌ Archivo 'PLAN de Capacitacion.pptx' no encontrado."

        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        texto += " ".join(run.text for run in paragraph.runs) + "\n"

        if not texto.strip():
            return "❌ El archivo PowerPoint está vacío o no contiene texto accesible."

    except Exception as e:
        print(f"❌ Error al cargar el archivo PowerPoint {path}: {e}")
        return f"❌ Error al cargar la documentación. Por favor, asegúrate de que el archivo '{os.path.basename(path)}' existe y es accesible."

    return texto.strip()


conocimiento_jira = cargar_conocimiento(PLAN_CAPACITACION_FILE)


class GeminiWorker(QRunnable):
    def __init__(self, question, knowledge):
        super().__init__()
        self.question = question
        self.knowledge = knowledge
        self.signals = WorkerSignals()

    def run(self):
        print("💡 Generando respuesta con Gemini...")
        prompt_parts = [
            {
                "text": "Eres un asistente de soporte técnico de Jira y también un Tester Senior con conocimiento en ISTQB."},
            {
                "text": "Responde a las preguntas de Jira utilizando la siguiente documentación, y responde a preguntas generales de testing con tus conocimientos de ISTQB."},
            {"text": f"Documentación de Jira: {self.knowledge}"},
            {"text": f"Pregunta del usuario: {self.question}"},
            {"text": "Respuesta clara y concisa:"}
        ]
        payload = {
            "contents": [
                {
                    "role": "user",
                    "parts": prompt_parts
                }
            ]
        }
        try:
            with httpx.Client(timeout=30.0) as client:
                response = client.post(
                    API_URL.format(API_KEY=API_KEY),
                    headers={"Content-Type": "application/json"},
                    json=payload,
                )
                response.raise_for_status()

                result = response.json()
                print("✅ Respuesta de Gemini recibida.")

                if result.get("candidates") and len(result["candidates"]) > 0 and \
                        result["candidates"][0].get("content") and \
                        result["candidates"][0]["content"].get("parts") and \
                        len(result["candidates"][0]["content"]["parts"]) > 0:
                    respuesta_llm = result["candidates"][0]["content"]["parts"][0].get("text", "").strip()
                    self.signals.result.emit(respuesta_llm)
                else:
                    error_msg = f"❌ Estructura de respuesta inesperada de Gemini: {result}"
                    self.signals.error.emit(error_msg)

        except httpx.RequestError as e:
            error_msg = f"❌ Error de red o conexión al llamar a la API de Gemini: {e}"
            self.signals.error.emit(error_msg)
        except httpx.HTTPStatusError as e:
            error_msg = f"❌ Error del servicio de IA (Código: {e.response.status_code}): {e.response.text}"
            self.signals.error.emit(error_msg)
        except json.JSONDecodeError as e:
            error_msg = f"❌ Error al decodificar la respuesta JSON: {e}"
            self.signals.error.emit(error_msg)
        except Exception as e:
            error_msg = f"❌ Error inesperado: {e}"
            self.signals.error.emit(error_msg)


class WorkerSignals(QObject):
    """Define las señales disponibles desde un hilo de trabajo."""
    result = Signal(str)
    error = Signal(str)


class JiraAssistantApp(QWidget):
    def __init__(self, stacked_widget):
        super().__init__()
        self.setWindowTitle("Asistente de Jira y QA")
        self.setGeometry(100, 100, 1000, 750)
        self.conocimiento_jira = conocimiento_jira
        self.stacked_widget = stacked_widget

        # --- Estilos Locales ---
        self.setStyleSheet("""
            QWidget {
                background-color: #1E1E2F;
                color: #F5F5F5;
                font-family: 'Segoe UI';
                font-size: 14px;
            }
            QFrame {
                background-color: #2A2A40;
                border-radius: 12px;
                border: 1px solid #3C3C5C;
            }
            QLabel {
                font-size: 16px;
                font-weight: bold;
                color: #B0C4DE;
            }
            QTextEdit, QListWidget {
                background-color: #3A3A50;
                border: 1px solid #5C5C7A;
                color: #F5F5F5;
                border-radius: 5px;
                padding: 5px;
            }
            QPushButton {
                background-color: #3F51B5;
                color: white;
                border-radius: 8px;
                padding: 10px;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: #5C6BC0;
            }
            QProgressBar {
                background-color: #2A2A40;
                border: 1px solid #3C3C5C;
                border-radius: 5px;
                text-align: center;
                color: #F5F5F5;
            }
            QProgressBar::chunk {
                background-color: #4CAF50;
                border-radius: 5px;
            }
        """)

        # Main layout
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(20, 20, 20, 20)

        # Botón de Home (casita)
        btn_home = QPushButton("🏠 Menú Principal")
        btn_home.setFixedSize(150, 40)
        btn_home.setStyleSheet(
            "QPushButton { background-color: #3F51B5; color: white; border-radius: 8px; font-size: 14px; } QPushButton:hover { background-color: #5C6BC0; }")
        btn_home.clicked.connect(lambda: self.stacked_widget.setCurrentIndex(0))

        lbl_appname = QLabel("🚀 Chat Assistant")
        lbl_appname.setStyleSheet("font-size: 18px; font-weight: bold; color: #FFD700;margin-right: 100px;")

        top_bar = QHBoxLayout()
        top_bar.addWidget(btn_home, alignment=Qt.AlignLeft)
        top_bar.addStretch()
        top_bar.addWidget(lbl_appname, alignment=Qt.AlignCenter)
        top_bar.addStretch()
        main_layout.addLayout(top_bar)

        # --- Contenedor superior para preguntas y sugerencias ---
        top_container = QHBoxLayout()
        top_container.setSpacing(20)

        # --- Card de Entrada de Pregunta ---
        input_card = QFrame()
        input_card.setFrameShape(QFrame.StyledPanel)
        input_layout = QVBoxLayout(input_card)

        lbl_info = QLabel("Realiza tu consulta:")
        lbl_info.setStyleSheet("font-size: 14px; font-weight: bold;")
        self.input_text = QTextEdit()
        self.input_text.setPlaceholderText("Ej: ¿Cómo crear un ticket en Jira? o ¿Qué es el smoke testing?")
        input_layout.addWidget(lbl_info)
        input_layout.addWidget(self.input_text)
        top_container.addWidget(input_card, 2)

        # --- Card de Preguntas Sugeridas ---
        sugg_card = QFrame()
        sugg_card.setFrameShape(QFrame.StyledPanel)
        sugg_layout = QVBoxLayout(sugg_card)

        lbl_sugg = QLabel("Preguntas sugeridas")
        lbl_sugg.setStyleSheet("font-size: 14px; font-weight: bold;")
        self.list_sugg = QListWidget()
        self.list_sugg.setSelectionMode(QListWidget.SingleSelection)
        self.list_sugg.setFixedWidth(250)

        questions = [
            "¿Cómo crear un ticket en Jira?",
            "¿Qué es un épico en Jira?",
            "¿Cómo se relaciona una historia de usuario con un épico?",
            "¿Qué es la matriz de trazabilidad?",
            "¿Qué es el 'Test Driven Development'?",
            "Explica los niveles de prueba según ISTQB."
        ]
        for q in questions:
            self.list_sugg.addItem(q)
        self.list_sugg.itemClicked.connect(self.select_suggested_question)

        sugg_layout.addWidget(lbl_sugg)
        sugg_layout.addWidget(self.list_sugg)
        top_container.addWidget(sugg_card, 1)

        main_layout.addLayout(top_container)

        # --- Contenedor de botones ---
        btn_container = QHBoxLayout()
        self.send_button = QPushButton("🚀 Enviar")
        self.send_button.clicked.connect(self.send_query)
        self.clear_button = QPushButton("🧹 Limpiar pantalla")
        self.clear_button.clicked.connect(self.clear_screen)

        btn_container.addWidget(self.send_button)
        btn_container.addWidget(self.clear_button)

        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        btn_container.addWidget(self.progress_bar)

        main_layout.addLayout(btn_container)

        # --- Card de Respuesta ---
        response_card = QFrame()
        response_card.setFrameShape(QFrame.StyledPanel)
        response_layout = QVBoxLayout(response_card)

        lbl_response = QLabel("Respuesta del modelo")
        lbl_response.setStyleSheet("font-size: 14px; font-weight: bold;")
        self.response_text = QTextEdit()
        self.response_text.setReadOnly(True)
        self.response_text.setPlaceholderText("La respuesta aparecerá aquí...")
        response_layout.addWidget(lbl_response)
        response_layout.addWidget(self.response_text)
        main_layout.addWidget(response_card)

        # Conectando la tecla Enter
        self.input_text.installEventFilter(self)

        # Pool de hilos
        self.threadpool = QThreadPool()
        # print(f"🚀 Iniciando con un máximo de {self.threadpool.maxThreadCount()} hilos disponibles.")

    def eventFilter(self, obj, event):
        if obj is self.input_text and event.type() == QEvent.KeyPress:
            if event.key() == Qt.Key_Return and not (event.modifiers() & Qt.ShiftModifier):
                self.send_query()
                return True
        return super().eventFilter(obj, event)

    def send_query(self):
        pregunta = self.input_text.toPlainText().strip()
        if not pregunta:
            QMessageBox.warning(self, "Alerta", "Por favor, escribe una pregunta.")
            return

        self.input_text.setEnabled(False)
        self.send_button.setEnabled(False)
        self.clear_button.setEnabled(False)
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)
        self.response_text.setPlaceholderText("Generando respuesta...")

        worker = GeminiWorker(pregunta, self.conocimiento_jira)
        worker.signals.result.connect(self.handle_response)
        worker.signals.error.connect(self.handle_error)
        self.threadpool.start(worker)

    def handle_response(self, response_text):
        self.response_text.setText(response_text)
        self.reset_ui()

    def handle_error(self, error_message):
        QMessageBox.critical(self, "Error", error_message)
        self.response_text.setText(f"Error: {error_message}")
        self.reset_ui()

    def reset_ui(self):
        self.input_text.setEnabled(True)
        self.send_button.setEnabled(True)
        self.clear_button.setEnabled(True)
        self.progress_bar.setVisible(False)
        self.progress_bar.setRange(0, 100)
        self.response_text.setPlaceholderText("La respuesta aparecerá aquí...")
        self.input_text.setFocus()

    def select_suggested_question(self, item):
        self.input_text.setText(item.text())

    def clear_screen(self):
        self.input_text.clear()
        self.response_text.clear()
        self.input_text.setFocus()